import os
import json
from pathlib import Path
from typing import Iterable

from datapizzai.clients.openai_client import OpenAIClient
from datapizzai.embedders.client_embedder import NodeEmbedder
from demo_minimo.components.text_splitter import TextSplitter
from datapizzai.type.type import Chunk


DATA_DIR = Path(__file__).parent / "data"
INDEX_DIR = Path(__file__).parent / "index"
INDEX_FILE = INDEX_DIR / "chunks.jsonl"
REMOTE_CACHE_DIR = Path(__file__).parent / "remote_cache"


SUPPORTED_EXTS = {".txt", ".md", ".pdf", ".pptx", ".py", ".ipynb", ".docx", ".html", ".htm", ".xlsx"}


def iter_files(root: Path) -> Iterable[Path]:
    for dirpath, dirnames, filenames in os.walk(root, followlinks=True):
        for fn in filenames:
            p = Path(dirpath) / fn
            if p.suffix.lower() in SUPPORTED_EXTS:
                yield p


def load_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF
        except Exception as e:
            raise RuntimeError("Install pymupdf to parse PDFs: pip install pymupdf") from e
        text_parts = []
        with fitz.open(str(path)) as doc:
            for page in doc:
                text_parts.append(page.get_text())
        return "\n".join(text_parts)
    if ext == ".pptx":
        try:
            from pptx import Presentation
        except Exception as e:
            raise RuntimeError("Install python-pptx to parse PPTX: pip install python-pptx") from e
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts)
    if ext == ".py":
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".ipynb":
        import json as _json

        try:
            nb = _json.loads(path.read_text(encoding="utf-8", errors="ignore"))
            cells = nb.get("cells", [])
            cell_texts = []
            for c in cells:
                src = c.get("source", [])
                if isinstance(src, list):
                    cell_texts.append("".join(src))
                elif isinstance(src, str):
                    cell_texts.append(src)
            return "\n".join(cell_texts)
        except Exception:
            return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".docx":
        try:
            import docx
        except Exception as e:
            raise RuntimeError("Install python-docx to parse DOCX: pip install python-docx") from e
        doc = docx.Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(parts)
    if ext in {".html", ".htm"}:
        try:
            from bs4 import BeautifulSoup
        except Exception as e:
            raise RuntimeError("Install beautifulsoup4 to parse HTML: pip install beautifulsoup4") from e
        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        # remove scripts/styles
        for s in soup(["script", "style"]):
            s.extract()
        text = soup.get_text(separator="\n")
        return "\n".join([line.strip() for line in text.splitlines() if line.strip()])
    if ext == ".xlsx":
        try:
            import openpyxl
        except Exception as e:
            raise RuntimeError("Install openpyxl to parse XLSX: pip install openpyxl") from e
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) if v is not None else "" for v in row]
                if any(vals):
                    parts.append("\t".join(vals))
        return "\n".join(parts)
    return ""


def save_chunks(chunks: list[Chunk]):
    INDEX_DIR.mkdir(parents=True, exist_ok=True)
    with open(INDEX_FILE, "w", encoding="utf-8") as f:
        for c in chunks:
            obj = {
                "id": c.id,
                "text": c.text,
                "metadata": c.metadata,
                "embedding": c.embeddings[0].vector if c.embeddings else [],
            }
            f.write(json.dumps(obj, ensure_ascii=False) + "\n")


def main(subdir: str | None = None, batch_size: int = 64, root_override: Path | None = None):
    api_key = os.getenv("OPENAI_API_KEY")
    embed_model = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-small")
    if not api_key:
        raise RuntimeError("Set OPENAI_API_KEY")
    root = root_override or (DATA_DIR / subdir if subdir else DATA_DIR)
    if not root.exists():
        raise RuntimeError(f"Data dir not found: {root}")

    client = OpenAIClient(api_key=api_key, model=embed_model)
    splitter = TextSplitter(max_char=1500, overlap=150)
    embedder = NodeEmbedder(client=client, model_name=embed_model, batch_size=batch_size)

    all_chunks: list[Chunk] = []
    for path in iter_files(root):
        text = load_text(path)
        chunks = splitter._run(text)
        for c in chunks:
            c.metadata.update({"source": str(path)})
        all_chunks.extend(chunks)

    all_chunks = embedder._run(all_chunks)
    save_chunks(all_chunks)
    print(f"Saved {len(all_chunks)} chunks to {INDEX_FILE}")


if __name__ == "__main__":
    # Allow optional subdir via env for quick runs
    subdir = os.getenv("HB_INGEST_SUBDIR")
    main(subdir=subdir)


def _extract_drive_id(url_or_id: str) -> str | None:
    url = url_or_id.strip().strip('"').strip("'")
    if not url:
        return None
    if "/folders/" in url:
        try:
            return url.split("/folders/")[-1].split("?")[0]
        except Exception:
            return None
    if "/drive/folders/" in url:
        try:
            return url.split("/drive/folders/")[-1].split("?")[0]
        except Exception:
            return None
    # Assume it's already an ID
    return url


def sync_gdrive(folder_urls: list[str], *, allowed_exts: set[str] | None = None) -> int:
    """Sync files from Google Drive folders into REMOTE_CACHE_DIR.

    Requires service account JSON path in env GOOGLE_DRIVE_SERVICE_ACCOUNT_JSON.
    The Drive folders must be shared with the service account email.
    """
    allowed_exts = allowed_exts or SUPPORTED_EXTS
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseDownload
        from googleapiclient.errors import HttpError
        import io, time
    except Exception as e:
        raise RuntimeError(
            "Servono le librerie Google API (google-api-python-client, google-auth). Installa e riprova."
        ) from e

    cred_path = os.getenv("GOOGLE_DRIVE_SERVICE_ACCOUNT_JSON")
    if not cred_path or not Path(cred_path).exists():
        raise RuntimeError(
            "Imposta GOOGLE_DRIVE_SERVICE_ACCOUNT_JSON con il path al JSON del service account e assicurati che le cartelle siano condivise con quell'account."
        )

    scopes = ["https://www.googleapis.com/auth/drive.readonly"]
    creds = service_account.Credentials.from_service_account_file(cred_path, scopes=scopes)
    service = build("drive", "v3", credentials=creds)

    REMOTE_CACHE_DIR.mkdir(parents=True, exist_ok=True)

    def _with_retry(fn, *, max_attempts: int = 5, base_delay: float = 0.8):
        for attempt in range(1, max_attempts + 1):
            try:
                return fn()
            except HttpError as e:  # type: ignore[name-defined]
                status = getattr(e, 'status_code', None) or getattr(getattr(e, 'resp', None), 'status', None)
                reason = ''
                try:
                    reason = (e.error_details or [{}])[0].get('reason', '')  # pyright: ignore[reportAttributeAccessIssue]
                except Exception:
                    pass
                if status in (500, 502, 503, 504) or 'internalError' in str(e) or 'rateLimitExceeded' in str(e):
                    time.sleep(base_delay * attempt)
                    continue
                raise
            except Exception:
                time.sleep(base_delay * attempt)
                if attempt == max_attempts:
                    raise

    def list_children(folder_id: str) -> list[dict]:
        files: list[dict] = []
        page_token: str | None = None
        print(f"Listing children of folder ID: {folder_id}")
        
        while True:
            def _do():
                return (
                    service.files()
                    .list(
                        q=f"'{folder_id}' in parents and trashed=false",
                        fields="nextPageToken, files(id, name, mimeType)",
                        pageToken=page_token,
                        supportsAllDrives=True,
                        includeItemsFromAllDrives=True,
                    )
                    .execute()
                )
                    
            try:
                res = _with_retry(_do)
                files.extend(res.get("files", []))
                page_token = res.get("nextPageToken")
                if not page_token:
                    break
            except Exception as e:
                print(f"Error listing folder {folder_id}: {e}")
                break
                
        print(f"Found {len(files)} files in folder {folder_id}")
        return files

    def download_file(file_id: str, target_path: Path) -> bool:
        # Skip if file already exists
        if target_path.exists():
            print(f"Skipping existing file: {target_path.name}")
            return False
        
        print(f"Downloading: {target_path.name}")
        def _start():
            request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
            fh = io.FileIO(target_path, "wb")
            return MediaIoBaseDownload(fh, request)
        
        try:
            downloader = _with_retry(_start)
            done = False
            while not done:
                try:
                    status, done = downloader.next_chunk()
                except HttpError as e:  # type: ignore[name-defined]
                    # retry by reinitializing downloader
                    downloader = _start()
                    continue
            print(f"Downloaded: {target_path.name}")
            return True
        except Exception as e:
            print(f"Error downloading {target_path.name}: {e}")
            # Remove partial file if exists
            if target_path.exists():
                target_path.unlink()
            return False

    def walk_and_download(folder_id: str, base_dir: Path) -> int:
        count = 0
        stack = [(folder_id, base_dir)]
        print(f"Starting download from folder ID: {folder_id}")
        while stack:
            fid, dir_path = stack.pop()
            dir_path.mkdir(parents=True, exist_ok=True)
            print(f"Processing folder: {dir_path.name}")
            children = list_children(fid)
            print(f"Found {len(children)} items in folder")
            for item in children:
                name = item["name"]
                mime = item.get("mimeType", "")
                if mime == "application/vnd.google-apps.folder":
                    print(f"Adding subfolder to queue: {name}")
                    stack.append((item["id"], dir_path / name))
                else:
                    ext = "." + name.split(".")[-1].lower() if "." in name else ""
                    if allowed_exts and ext not in allowed_exts:
                        print(f"Skipping unsupported file: {name}")
                        continue
                    target = dir_path / name
                    if download_file(item["id"], target):
                        count += 1
        print(f"Download completed. Total files downloaded: {count}")
        return count

    total = 0
    print(f"Starting sync_gdrive with {len(folder_urls)} URLs")
    for i, u in enumerate(folder_urls):
        print(f"Processing URL {i+1}/{len(folder_urls)}: {u}")
        fid = _extract_drive_id(u)
        if not fid:
            print(f"Could not extract folder ID from URL: {u}")
            continue
        print(f"Extracted folder ID: {fid}")
        downloaded = walk_and_download(fid, REMOTE_CACHE_DIR)
        total += downloaded
        print(f"Downloaded {downloaded} files from this URL")

    print(f"sync_gdrive completed. Total files downloaded: {total}")
    return total


